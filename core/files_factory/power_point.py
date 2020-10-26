from pptx import Presentation
from core.files_factory.base_file_type import FileFactory

import logging
logger = logging.getLogger(__name__)


class PowerPoint(FileFactory):
    
    def search(self, file_path, search_term, **kwargs):
        try:
            presentation = Presentation(str(file_path)) 
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape.text = shape.text.lower()
                        if search_term.lower() in shape.text:
                            return True
        except Exception as e:
            logger.error(f'Error while extract power point {file_path}: {e}')
        return False
